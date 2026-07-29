"""
File processing utilities for My Revision Helper.

This module handles:
- Image compression for OpenAI Vision API
- Text extraction from images (OCR via OpenAI)
- Text extraction from PDFs, including scans with no text layer
- Text extraction from PowerPoint presentations
- File processing orchestration
"""

from __future__ import annotations

import asyncio
import base64
import logging
from io import BytesIO
from typing import Any, Dict, List, Optional

from fastapi import UploadFile
from PIL import Image

# Set up logging
logger = logging.getLogger(__name__)

# Maximum file size for upload: 50MB (we'll compress/process as needed)
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB in bytes

# Maximum size for OpenAI Vision API (after base64 encoding, ~20MB raw = ~27MB base64)
# We'll compress images to stay under this limit
MAX_OPENAI_IMAGE_SIZE = 20 * 1024 * 1024  # 20MB raw image (before base64)


def compress_image(image_bytes: bytes, max_size: int = MAX_OPENAI_IMAGE_SIZE, quality: int = 85) -> bytes:
    """
    Compress an image to fit within size limits while preserving quality.
    
    Args:
        image_bytes: Original image bytes
        max_size: Maximum size in bytes (default: MAX_OPENAI_IMAGE_SIZE)
        quality: JPEG quality (1-100, default: 85)
        
    Returns:
        Compressed image bytes (JPEG format)
    """
    try:
        # Open image
        img = Image.open(BytesIO(image_bytes))
        
        # Convert RGBA to RGB if needed (for JPEG compatibility)
        if img.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
            img = background
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        
        # If already small enough, return as-is
        output = BytesIO()
        img.save(output, format='JPEG', quality=quality, optimize=True)
        compressed = output.getvalue()
        
        if len(compressed) <= max_size:
            return compressed
        
        # Need to resize - calculate scale factor
        original_size = len(image_bytes)
        scale_factor = (max_size / original_size) ** 0.5  # Square root for 2D scaling
        
        # Resize image
        new_width = int(img.width * scale_factor)
        new_height = int(img.height * scale_factor)
        
        # Ensure minimum dimensions for readability
        new_width = max(new_width, 800)
        new_height = max(new_height, 600)
        
        img_resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
        
        # Try different quality levels
        for q in range(quality, 50, -10):
            output = BytesIO()
            img_resized.save(output, format='JPEG', quality=q, optimize=True)
            compressed = output.getvalue()
            if len(compressed) <= max_size:
                logger.info(f"Compressed image: {len(image_bytes) / (1024*1024):.1f}MB -> {len(compressed) / (1024*1024):.1f}MB (quality={q})")
                return compressed
        
        # If still too large, use minimum quality
        output = BytesIO()
        img_resized.save(output, format='JPEG', quality=50, optimize=True)
        compressed = output.getvalue()
        logger.info(f"Compressed image to minimum: {len(image_bytes) / (1024*1024):.1f}MB -> {len(compressed) / (1024*1024):.1f}MB")
        return compressed
        
    except Exception as e:
        logger.error(f"Failed to compress image: {e}", exc_info=True)
        # Return original if compression fails
        return image_bytes


# A page with less real text than this is treated as a picture of a page rather
# than a page of text. Scanner apps produce PDFs with no text layer at all, and
# a page carrying only a figure caption is not worth trying to read either.
MIN_CHARS_PER_TEXT_PAGE = 40

# Scanned pages cost an OpenAI Vision call each, so a very long document is read
# up to this point rather than running up an unbounded bill on a mis-upload.
MAX_OCR_PAGES = 30

# How many pages to read at once. Sequential OCR of a 15-page scan would keep a
# child waiting a couple of minutes.
OCR_CONCURRENCY = 4

# 2x gives roughly 144 dpi, which is enough for handwriting without producing
# images so large they need compressing again.
PDF_RENDER_SCALE = 2.0


def _render_pdf_page(contents: bytes, page_index: int, scale: float = PDF_RENDER_SCALE) -> Optional[bytes]:
    """
    Render one PDF page to JPEG bytes.

    Uses pypdfium2, which arrives with pdfplumber and needs no system packages,
    so scanned PDFs work on a plain Python image.
    """
    try:
        import pypdfium2

        document = pypdfium2.PdfDocument(BytesIO(contents))
        try:
            page = document[page_index]
            image = page.render(scale=scale).to_pil()
            if image.mode != "RGB":
                image = image.convert("RGB")
            buffer = BytesIO()
            image.save(buffer, format="JPEG", quality=80, optimize=True)
            return buffer.getvalue()
        finally:
            document.close()
    except ImportError:
        logger.error("pypdfium2 not installed - cannot read scanned PDFs")
        return None
    except Exception as e:
        logger.warning(f"Could not render PDF page {page_index + 1}: {e}")
        return None


async def _ocr_pdf_pages(
    contents: bytes, page_numbers: List[int], client: Any, filename: str
) -> Dict[int, str]:
    """
    Read scanned pages with OpenAI Vision, several at a time.

    Returns page number -> text for the pages that could be read.
    """
    if not page_numbers or client is None:
        return {}

    if len(page_numbers) > MAX_OCR_PAGES:
        logger.warning(
            f"{filename} has {len(page_numbers)} scanned pages; reading the first {MAX_OCR_PAGES}"
        )
        page_numbers = page_numbers[:MAX_OCR_PAGES]

    limit = asyncio.Semaphore(OCR_CONCURRENCY)

    async def read(page_num: int):
        async with limit:
            rendered = await asyncio.to_thread(_render_pdf_page, contents, page_num - 1)
            if not rendered:
                return page_num, None
            text = await asyncio.to_thread(
                _vision_ocr, rendered, client, f"{filename} page {page_num}"
            )
            return page_num, text

    logger.info(f"Reading {len(page_numbers)} scanned page(s) of {filename} with Vision")
    results = await asyncio.gather(*(read(n) for n in page_numbers), return_exceptions=True)

    out: Dict[int, str] = {}
    for result in results:
        if isinstance(result, Exception):
            logger.warning(f"A page of {filename} could not be read: {result}")
            continue
        page_num, text = result
        if text and text.strip():
            out[page_num] = text.strip()

    return out


async def extract_text_from_pdf(file: UploadFile, client: Any = None) -> Optional[str]:
    """
    Extract text from a PDF file.

    Handles both kinds of PDF that turn up. One is a real document with a text
    layer, read directly. The other is a scan or a phone photo — anything from a
    scanner app is a picture of a page with no text in it at all — which is read
    with OpenAI Vision, page by page, when a client is available.

    Args:
        file: Uploaded PDF file
        client: OpenAI client, needed only to read scanned pages

    Returns:
        Extracted text, or None if extraction fails
    """
    try:
        import pdfplumber
        
        # Read file content
        contents = await file.read()
        
        # Check file size - PDFs can be large, but we process them locally
        if len(contents) > MAX_FILE_SIZE:
            logger.warning(f"PDF {file.filename} is large ({len(contents) / (1024*1024):.1f}MB), processing in chunks...")
        
        # Use pdfplumber to extract text (processes all pages)
        page_text_by_number: Dict[int, str] = {}
        scanned_pages: List[int] = []

        with pdfplumber.open(BytesIO(contents)) as pdf:
            total_pages = len(pdf.pages)
            logger.info(f"Processing PDF {file.filename} with {total_pages} pages")
            
            for page_num, page in enumerate(pdf.pages, 1):
                try:
                    page_text = page.extract_text()
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num} of {file.filename}: {e}")
                    page_text = None

                if page_text and len(page_text.strip()) >= MIN_CHARS_PER_TEXT_PAGE:
                    page_text_by_number[page_num] = page_text.strip()
                    continue

                # Not enough text to be a text page. If there is something drawn
                # on it, it is a scan and worth reading properly; a genuinely
                # blank page is skipped.
                if page.images or page_text:
                    scanned_pages.append(page_num)
                    if page_text and page_text.strip():
                        page_text_by_number[page_num] = page_text.strip()

        if scanned_pages and client is not None:
            for page_num, text in (
                await _ocr_pdf_pages(contents, scanned_pages, client, file.filename or "PDF")
            ).items():
                # OCR of a scan beats the stray characters the text layer had.
                page_text_by_number[page_num] = text
        elif scanned_pages:
            logger.warning(
                f"{file.filename} has {len(scanned_pages)} page(s) with no text layer "
                "and no OpenAI client was available to read them"
            )

        text_parts = [
            f"--- Page {num} ---\n{page_text_by_number[num]}"
            for num in sorted(page_text_by_number)
        ]

        extracted_text = "\n\n".join(text_parts)
        logger.info(
            f"Extracted {len(extracted_text)} characters from PDF {file.filename} "
            f"({total_pages} pages, {len(scanned_pages)} scanned)"
        )
        return extracted_text if extracted_text else None
            
    except ImportError:
        logger.error("pdfplumber not installed - cannot extract text from PDFs")
        return None
    except Exception as e:
        logger.error(f"Failed to extract text from PDF {file.filename}: {e}", exc_info=True)
        return None


async def extract_text_from_pptx(file: UploadFile) -> Optional[str]:
    """
    Extract text from a PowerPoint presentation file.
    
    Args:
        file: Uploaded PPTX file
        
    Returns:
        Extracted text, or None if extraction fails
    """
    try:
        from pptx import Presentation
        
        # Read file content
        contents = await file.read()
        
        # Check file size - PPTX can be large, but we process them locally
        if len(contents) > MAX_FILE_SIZE:
            logger.warning(f"PPTX {file.filename} is large ({len(contents) / (1024*1024):.1f}MB), processing all slides...")
        
        # Use python-pptx to extract text (processes all slides)
        prs = Presentation(BytesIO(contents))
        text_parts = []
        total_slides = len(prs.slides)
        logger.info(f"Processing PPTX {file.filename} with {total_slides} slides")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)
            
            if slide_texts:
                slide_text = f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts)
                text_parts.append(slide_text)
        
        extracted_text = "\n\n".join(text_parts)
        logger.info(f"Extracted {len(extracted_text)} characters from PPTX {file.filename} ({total_slides} slides)")
        return extracted_text if extracted_text else None
        
    except ImportError:
        logger.error("python-pptx not installed - cannot extract text from PowerPoint files")
        return None
    except Exception as e:
        logger.error(f"Failed to extract text from PPTX {file.filename}: {e}", exc_info=True)
        return None


def _docx_text_from_bytes(contents: bytes) -> Optional[str]:
    """
    Pull text out of a .docx without requiring python-docx.

    A .docx is a zip of XML. Paragraphs and table cells are converted to lines
    and pipe-separated columns respectively, which keeps the shape of the
    tables that study plans and trackers rely on.
    """
    import html
    import re
    import zipfile

    with zipfile.ZipFile(BytesIO(contents)) as archive:
        if "word/document.xml" not in archive.namelist():
            return None
        xml = archive.read("word/document.xml").decode("utf-8", errors="replace")

    # Preserve document structure before stripping tags.
    xml = xml.replace("</w:tab>", "\t")
    xml = xml.replace("<w:br/>", "\n").replace("<w:br />", "\n")
    xml = xml.replace("</w:tc>", " | ")
    xml = xml.replace("</w:tr>", "\n")
    xml = xml.replace("</w:p>", "\n")

    text = html.unescape(re.sub(r"<[^>]+>", "", xml))

    lines: List[str] = []
    for raw_line in text.split("\n"):
        line = raw_line.rstrip().rstrip(" |")
        if not line.strip():
            # Collapse runs of blank lines to a single separator.
            if lines and lines[-1] != "":
                lines.append("")
            continue
        lines.append(line)

    result = "\n".join(lines).strip()
    return result or None


def _xlsx_text_from_bytes(contents: bytes) -> Optional[str]:
    """
    Pull text out of an .xlsx without requiring openpyxl.

    Each sheet becomes a titled section of pipe-separated rows, so trackers and
    score logs stay readable to both a human and the model.
    """
    import html
    import re
    import zipfile

    with zipfile.ZipFile(BytesIO(contents)) as archive:
        names = archive.namelist()
        if "xl/workbook.xml" not in names:
            return None

        # Shared strings are referenced by index from the cells.
        shared: List[str] = []
        if "xl/sharedStrings.xml" in names:
            shared_xml = archive.read("xl/sharedStrings.xml").decode("utf-8", errors="replace")
            for item in re.findall(r"<si>(.*?)</si>", shared_xml, re.S):
                parts = re.findall(r"<t[^>]*>(.*?)</t>", item, re.S)
                shared.append(html.unescape(re.sub(r"<[^>]+>", "", "".join(parts))))

        workbook_xml = archive.read("xl/workbook.xml").decode("utf-8", errors="replace")
        sheets = re.findall(r'<sheet[^>]*name="([^"]+)"[^>]*r:id="(rId\d+)"', workbook_xml)

        rel_map: Dict[str, str] = {}
        if "xl/_rels/workbook.xml.rels" in names:
            rels_xml = archive.read("xl/_rels/workbook.xml.rels").decode("utf-8", errors="replace")
            rel_map = dict(re.findall(r'Id="(rId\d+)"[^>]*Target="([^"]+)"', rels_xml))

        sections: List[str] = []
        for sheet_name, rel_id in sheets:
            target = rel_map.get(rel_id, "")
            if not target:
                continue
            path = "xl/" + target.lstrip("/").removeprefix("xl/")
            if path not in names:
                continue

            sheet_xml = archive.read(path).decode("utf-8", errors="replace")
            rows: List[str] = []
            for row_xml in re.findall(r"<row[^>]*>(.*?)</row>", sheet_xml, re.S):
                values: List[str] = []
                for attrs, body, self_closing_attrs in re.findall(
                    r"<c\b([^>]*)>(.*?)</c>|<c\b([^>]*)/>", row_xml, re.S
                ):
                    attrs = attrs or self_closing_attrs
                    inline = re.findall(r"<t[^>]*>(.*?)</t>", body, re.S)
                    if inline:
                        value = "".join(inline)
                    else:
                        v_match = re.search(r"<v>(.*?)</v>", body, re.S)
                        if not v_match:
                            values.append("")
                            continue
                        value = v_match.group(1)
                        type_match = re.search(r't="([^"]+)"', attrs)
                        if type_match and type_match.group(1) == "s":
                            index = int(value)
                            value = shared[index] if index < len(shared) else value
                    values.append(html.unescape(str(value)))

                line = " | ".join(values).rstrip(" |")
                if line.strip():
                    rows.append(line)

            if rows:
                sections.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))

    result = "\n\n".join(sections).strip()
    return result or None


async def extract_text_from_docx(file: UploadFile) -> Optional[str]:
    """
    Extract text from a Word document, preserving paragraphs and table layout.

    Args:
        file: Uploaded .docx file

    Returns:
        Extracted text, or None if extraction fails
    """
    try:
        contents = await file.read()
        text = _docx_text_from_bytes(contents)
        if text:
            logger.info(f"Extracted {len(text)} characters from DOCX {file.filename}")
        return text
    except Exception as e:
        logger.error(f"Failed to extract text from DOCX {file.filename}: {e}", exc_info=True)
        return None


async def extract_text_from_xlsx(file: UploadFile) -> Optional[str]:
    """
    Extract text from a spreadsheet, one titled section per sheet.

    Args:
        file: Uploaded .xlsx file

    Returns:
        Extracted text, or None if extraction fails
    """
    try:
        contents = await file.read()
        text = _xlsx_text_from_bytes(contents)
        if text:
            logger.info(f"Extracted {len(text)} characters from XLSX {file.filename}")
        return text
    except Exception as e:
        logger.error(f"Failed to extract text from XLSX {file.filename}: {e}", exc_info=True)
        return None


def _strip_code_fence(text: Optional[str]) -> Optional[str]:
    """
    Remove a ``` wrapper the model sometimes puts around a whole page.

    Left in, it becomes noise in the marking prompt and can swallow the first
    line of a page.
    """
    if not text:
        return text

    stripped = text.strip()
    if not stripped.startswith("```"):
        return text

    lines = stripped.split("\n")
    # Drop the opening fence, along with any language tag on it.
    lines = lines[1:]
    if lines and lines[-1].strip().startswith("```"):
        lines = lines[:-1]
    return "\n".join(lines).strip()


OCR_SYSTEM_PROMPT = """You are transcribing a page of school work so it can be marked. The page is \
usually a scan or photo of an exam paper a student has written on by hand.

Transcribe the page exactly as it is. Do not solve anything, correct anything, or tidy \
up the student's mistakes: a wrong answer must come through wrong, because it is being \
marked.

SEPARATE THE PRINTED PAPER FROM THE STUDENT'S WRITING
This matters more than anything else on this page. Whoever marks this has to know which \
words are the question and which are the student's answer.
- Transcribe printed text normally.
- Put [written] at the start of every line, or part of a line, the student wrote by hand.
- On an answer line, keep the printed label: Answer: [written] 490  ......... (2)
- If the student wrote nothing for a question, write [no answer].

MATHS
- Write maths in plain text with real symbols: × ÷ ± ≤ ≥ ≠ ≈ √ π ° ² ³ ½ ¾ ⅓.
- Powers as x² where you can, x^7 only when the exponent will not fit as a superscript.
- Fractions as 3/4, mixed numbers as 2 1/2, and a/b for algebraic fractions.
- Never use LaTeX, dollar signs, \\frac, \\times, align blocks or markdown code fences.
- Keep every line of the student's working on its own line, in the order it was written, \
including crossings out (write [crossed out] before them). Working earns method marks, \
so none of it may be dropped or merged.
- Copy digits exactly as written even if the arithmetic is wrong.

FIGURES, GRAPHS, CHARTS AND DIAGRAMS
Never skip a figure and never reduce one to a bare name: on many papers the drawing IS \
the answer and is worth several marks. For each one write a single block:

[FIGURE: printed | drawn by student — what it is. Then everything that could be marked.]

Say what can actually be measured or checked:
- Pie chart: every sector, its label, and its size as an angle or fraction of the circle.
- Bar chart or histogram: each bar's label and its height read off the scale.
- Line or scatter graph: the axis labels and ranges, the points plotted with their \
coordinates, and whether a line or curve has been drawn through them.
- Shape or construction: labelled lengths, angles, right-angle and equal-length marks.
- Table: transcribe it as rows of text, not as a figure.
State whether it is printed on the paper or drawn by the student. If something is too \
faint to read, say so rather than guessing at it.

If any part is illegible write [illegible]. Return only the transcription, with no \
commentary, preamble or explanation."""


def _vision_ocr(contents: bytes, client: Any, label: str) -> Optional[str]:
    """
    Read the text in one image with OpenAI Vision.

    Synchronous and takes bytes rather than an upload, so it serves both a
    directly uploaded photo and a page rendered out of a scanned PDF.
    """
    try:
        original_size = len(contents)

        # Check file size - compress if too large for OpenAI
        if original_size > MAX_OPENAI_IMAGE_SIZE:
            logger.info(f"Image {label} is large ({original_size / (1024*1024):.1f}MB), compressing...")
            try:
                contents = compress_image(contents, max_size=MAX_OPENAI_IMAGE_SIZE)
                logger.info(f"Compressed {label}: {original_size / (1024*1024):.1f}MB -> {len(contents) / (1024*1024):.1f}MB")
            except ImportError:
                logger.warning("Pillow not available - cannot compress image. Install Pillow for large image support.")
                logger.error(f"Image {label} too large ({original_size / (1024*1024):.1f}MB) and compression unavailable")
                return None
            except Exception as e:
                logger.error(f"Failed to compress image {label}: {e}", exc_info=True)
                return None

        base64_image = base64.b64encode(contents).decode('utf-8')

        # After compression, images are always JPEG
        # (compression converts all formats to JPEG for consistency and size)
        image_format = "image/jpeg"

        # Reading a page of a child's handwriting is one of the harder things
        # asked of a model here, and a misread digit becomes a lost mark, so
        # this uses the better model rather than the everyday one.
        from .llm import chat_completion, get_reasoning_model

        response = chat_completion(
            client,
            model=get_reasoning_model(),
            messages=[
                {"role": "system", "content": OCR_SYSTEM_PROMPT},
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Extract all text from this image. Preserve the structure and formatting as much as possible.",
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{image_format};base64,{base64_image}",
                            },
                        },
                    ],
                },
            ],
            max_tokens=2000,
        )

        extracted_text = _strip_code_fence(response.choices[0].message.content)
        logger.info(f"Extracted {len(extracted_text or '')} characters from {label}")
        return extracted_text

    except Exception as e:
        logger.error(f"Failed to extract text from {label}: {e}", exc_info=True)
        return None


async def extract_text_from_image(file: UploadFile, client: Any) -> Optional[str]:
    """
    Extract text from an uploaded image using OpenAI Vision API.
    
    Args:
        file: Uploaded file (should be an image)
        client: OpenAI client instance
        
    Returns:
        Extracted text, or None if extraction fails
    """
    try:
        contents = await file.read()
    except Exception as e:
        logger.error(f"Failed to read image {file.filename}: {e}", exc_info=True)
        return None

    return _vision_ocr(contents, client, f"image {file.filename}")


async def process_uploaded_files(files: List[UploadFile], openai_client: Any) -> Dict[str, str]:
    """
    Process uploaded files and extract text from images, PDFs, and PowerPoint files.
    
    Args:
        files: List of uploaded files
        openai_client: OpenAI client instance (for image OCR)
        
    Returns:
        Dictionary mapping filename to extracted text
    """
    if not files:
        return {}
    
    extracted_texts = {}
    skipped_files = []
    
    for file in files:
        content_type = file.content_type or ""
        filename = file.filename or "unknown"
        filename_lower = filename.lower()
        
        # Reset file pointer (in case it was read before)
        await file.seek(0)
        
        try:
            # Handle PDF files
            if "pdf" in content_type or filename_lower.endswith('.pdf'):
                logger.info(f"Processing PDF file: {filename}")
                # The client is passed so a scanned PDF can be read as images.
                text = await extract_text_from_pdf(file, openai_client)
                if text:
                    extracted_texts[filename] = text
                else:
                    skipped_files.append(f"{filename} (PDF extraction failed or file too large)")
            
            # Handle PowerPoint files
            elif ("presentation" in content_type or "powerpoint" in content_type or 
                  filename_lower.endswith(('.pptx', '.ppt'))):
                logger.info(f"Processing PowerPoint file: {filename}")
                text = await extract_text_from_pptx(file)
                if text:
                    extracted_texts[filename] = text
                else:
                    skipped_files.append(f"{filename} (PowerPoint extraction failed or file too large)")
            
            # Handle Word documents
            elif ("wordprocessingml" in content_type or filename_lower.endswith('.docx')):
                logger.info(f"Processing Word document: {filename}")
                text = await extract_text_from_docx(file)
                if text:
                    extracted_texts[filename] = text
                else:
                    skipped_files.append(f"{filename} (Word extraction failed)")

            # Handle Excel spreadsheets
            elif ("spreadsheetml" in content_type or filename_lower.endswith('.xlsx')):
                logger.info(f"Processing spreadsheet: {filename}")
                text = await extract_text_from_xlsx(file)
                if text:
                    extracted_texts[filename] = text
                else:
                    skipped_files.append(f"{filename} (spreadsheet extraction failed)")

            # Handle plain text and markdown
            elif ("text/plain" in content_type or "markdown" in content_type or
                  filename_lower.endswith(('.txt', '.md'))):
                logger.info(f"Processing text file: {filename}")
                raw = await file.read()
                text = raw.decode("utf-8", errors="replace").strip()
                if text:
                    extracted_texts[filename] = text
                else:
                    skipped_files.append(f"{filename} (empty text file)")

            # Handle image files (requires OpenAI client)
            elif any(img_type in content_type for img_type in ["image/jpeg", "image/jpg", "image/png", "image/gif", "image/webp"]):
                if not openai_client:
                    skipped_files.append(f"{filename} (OpenAI client not available for image OCR)")
                    logger.warning(f"OpenAI client not available - cannot extract text from image {filename}")
                else:
                    logger.info(f"Processing image file: {filename}")
                    text = await extract_text_from_image(file, openai_client)
                    if text:
                        extracted_texts[filename] = text
                    else:
                        skipped_files.append(f"{filename} (image extraction failed or file too large)")
            
            else:
                skipped_files.append(f"{filename} (unsupported file type: {content_type})")
                logger.warning(f"Skipping unsupported file: {filename} (type: {content_type})")
        
        except Exception as e:
            logger.error(f"Error processing file {filename}: {e}", exc_info=True)
            skipped_files.append(f"{filename} (processing error: {str(e)[:50]})")
    
    if skipped_files:
        logger.info(f"Skipped {len(skipped_files)} file(s): {', '.join(skipped_files)}")
    
    if extracted_texts:
        logger.info(f"Successfully extracted text from {len(extracted_texts)} file(s)")
    
    return extracted_texts


